from tkinter import filedialog
from tkinter import *
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt
import os
from pptx.enum.text import PP_ALIGN
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.svm import SVC
from sklearn.model_selection import train_test_split
from sklearn.metrics import accuracy_score

def create_starting_slide(prs, title_text, subtitle_text, presenter_names):
    # Add a slide
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Check if the slide has a title placeholder
    title_shape = slide.shapes.title
    if not title_shape:
        # If no title placeholder, create one
        title_shape = slide.placeholders[0]

    # Remove any existing text boxes
    remove_existing_text_boxes(slide)

    # Add title text
    left = Inches(3)
    top = Inches(3)
    width = Inches(9)
    height = Inches(1)
    subtitle_shape = slide.shapes.add_textbox(left, top, width, height)
    subtitle_shape.text_frame.text = f"Title: {title_text}"
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(40)
    subtitle_shape.text_frame.paragraphs[0].font.bold = True

    # Add subtitle text
    left = Inches(3)
    top = Inches(4)
    width = Inches(9)
    height = Inches(1)
    subtitle_shape = slide.shapes.add_textbox(left, top, width, height)
    subtitle_shape.text_frame.text = f"Subtitle: {subtitle_text}"
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.bold = True

    # Add presented names
    left = Inches(10.5)
    top = Inches(5.5)
    width = Inches(9)
    height = Inches(1)
    presenter_shape = slide.shapes.add_textbox(left, top, width, height)
    presenter_shape.text_frame.text = "Presented By:"
    presenter_shape.text_frame.paragraphs[0].font.size = Pt(20)
    presenter_shape.text_frame.paragraphs[0].font.bold = True
    for name in presenter_names.split(","):
        p = presenter_shape.text_frame.add_paragraph()
        p.text = name.strip()
        p.font.size = Pt(16)

def create_table_slide(prs, excel_path, table_heading, rows, cols):
    # Check if the number of rows and columns exceed the slide limit
    if rows > 10 or cols > 10:
        print("The number of rows or columns exceeds the slide limit (10x10). Cannot create table.")
        return

    # Add a slide
    slide_layout = prs.slide_layouts[5]  # You may need to adjust the index based on your template
    slide = prs.slides.add_slide(slide_layout)

    # Insert table heading
    title_shape = slide.shapes.title
    title_shape.text = table_heading

    # Insert a table
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)
    gap_width = Inches(0.25)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Read data from Excel file
    wb = load_workbook(excel_path)
    ws = wb.active
    data = []
    for row in ws.iter_rows():
        data.append([cell.value for cell in row])

    # Populate the table with Excel data and set text wrapping
    for i in range(rows):
        for j in range(cols):
            if i < len(data) and j < len(data[i]):
                cell_text = str(data[i][j])
            else:
                cell_text = "None"  # Show "None" for empty cells
            cell = table.cell(i, j)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(12)  # Adjust font size as needed
            cell.text_frame.word_wrap = True

def add_thank_you_slide(prs):
    # Add a slide
    slide_layout = prs.slide_layouts[5]  # You may need to adjust the index based on your template
    slide = prs.slides.add_slide(slide_layout)

    # Remove any existing text boxes
    remove_existing_text_boxes(slide)

    # Add a text box with a smile emoji
    left = Inches(2)
    top = Inches(3.5)
    width = Inches(6)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Thank you for using SlideAI! 😊"
    p.font.size = Pt(30)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

def remove_existing_text_boxes(slide):
    # Remove any existing text boxes from the slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            slide.shapes._spTree.remove(shape._element)

def perform_sentiment_analysis(text_data, labels):
    # 1. Data Preparation
    # Assuming text_data is a list of text samples and labels is a list of corresponding sentiment labels

    # 2. Feature Extraction
    vectorizer = TfidfVectorizer(max_features=1000)  # Extract top 1000 features
    X = vectorizer.fit_transform(text_data)  # Convert text data into numerical feature vectors

    # 3. Model Training
    X_train, X_test, y_train, y_test = train_test_split(X, labels, test_size=0.2, random_state=42)  # Split data into train and test sets
    svm_model = SVC(kernel='linear')  # Initialize SVM model
    svm_model.fit(X_train, y_train)  # Train SVM model

    # 4. Prediction
    y_pred = svm_model.predict(X_test)  # Predict labels for test data

    # Evaluate model performance
    accuracy = accuracy_score(y_test, y_pred)
    print("Accuracy:", accuracy)

# Prompt user to add a PowerPoint template
print("Please select a PowerPoint template (only .pptx files)")
root = Tk()
root.withdraw()
template_path = filedialog.askopenfilename(initialdir="/", title="Select PowerPoint template file", filetypes=(("PowerPoint files", "*.pptx"), ("all files", "*.*")))

if not template_path:
    print("No PowerPoint template selected. Process terminated.")
    exit()

# Create a PowerPoint presentation
prs = Presentation(template_path)

# Prompt user to add title, subtitle, and presenter names for the starting slide
title_text = input("Enter title for the starting slide: ")
subtitle_text = input("Enter subtitle for the starting slide: ")
presenter_names = input("Enter presenter names (comma-separated): ")

# Create starting slide
create_starting_slide(prs, title_text, subtitle_text, presenter_names)

# Prompt user to enter the number of slides to create
num_slides = int(input("How many slides with tables do you want to create? (Max 5): "))
if num_slides > 5:
    print("The number of slides exceeds the limit (5). Cannot create more slides.")
    exit()

for i in range(num_slides):
    # Prompt user to add an Excel file
    print(f"Please select Excel file for slide {i+1} (only .xlsx files)")
    root = Tk()
    root.withdraw()
    excel_path = filedialog.askopenfilename(initialdir="/", title=f"Select Excel file for slide {i+1}", filetypes=(("Excel files", "*.xlsx"), ("all files", "*.*")))

    if not excel_path:
        print("No Excel file selected. Process terminated.")
        exit()

    # Prompt user to enter heading for the table slide
    table_heading = input(f"Enter heading for the table slide {i+1}: ")

    # Prompt user to enter the number of rows and columns for the table
    rows = int(input(f"Enter the number of rows for the table {i+1}: "))
    cols = int(input(f"Enter the number of columns for the table {i+1}: "))

    create_table_slide(prs, excel_path, table_heading, rows, cols)

# Add thank you slide
add_thank_you_slide(prs)

# Perform sentiment analysis (example)
text_data = ["This movie is great!", "I didn't like the food.", "The product quality is excellent."]
labels = ["Positive", "Negative", "Positive"]
perform_sentiment_analysis(text_data, labels)

# Prompt user to enter the name for the PowerPoint file
output_name = input("Enter the name for the PowerPoint file to be created (without extension): ")
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), f"{output_name}.pptx")

# Check if the output file already exists
if os.path.exists(output_path):
    choice = input("File already exists. Do you want to override it? (yes/no): ").lower()
    if choice != "yes":
        print("Process terminated.")
        exit()

# Save the PowerPoint presentation
prs.save(output_path)
print(f"Presentation saved to {output_path}")
